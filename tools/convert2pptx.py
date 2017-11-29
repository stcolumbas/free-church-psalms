import os
import re
import tempfile
import warnings
import zipfile

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from utils import (convert2pdf, load_scottish_psalter, load_sing_psalms,
                   make_output_folder, remove_folder, remove_markup,
                   zip_folder)

re_underlined = re.compile('\<underline\>.*\<\/underline\>')


def superscript_pptx(tmpfile):
    with zipfile.ZipFile(tmpfile, 'a') as f:
        slides = set([x for x in f.namelist() if x.startswith('ppt/slides/slide')])
        slides.remove('ppt/slides/slide1.xml')
        slides = list(slides)
        for slide in slides:
            text = f.open(slide).read().decode()
            text = re.sub(
                '<a:r><a:t>(\d+-*\d*)',
                '<a:r><a:rPr baseline="30000" dirty="0"/><a:t>\g<1></a:t></a:r><a:r><a:rPr dirty="0"/><a:t>',
                text
            )
            text = re.sub(
                '\n(\d+-*\d*)',
                '\n</a:t></a:r><a:r><a:rPr baseline="30000" dirty="0"/><a:t>\g<1></a:t></a:r><a:r><a:rPr dirty="0"/><a:t>',
                text
            )
            f.writestr(slide, text)


def add_underlined_text(p, text):
    match = re_underlined.search(text)
    if match is None:
        # no matches, just add the text as is
        r = p.add_run()
        r.text = text
        return
    # there is a match, extract the text
    span = match.span()
    beg = text[0:span[0]]
    underlined = text[span[0]:span[1]]
    underlined = underlined.replace('<underline>', '').replace('</underline>', '')
    rest = text[span[1]:]
    # add to slide:
    if beg:
        r = p.add_run()
        r.text = beg

    r = p.add_run()
    r.text = underlined
    r.font.underline = True

    # add rest of text:
    if rest:
        add_underlined_text(p, rest)


def add_slide(prs, v, ratio, underline):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = Inches(0)
    width = Inches(20)
    if ratio == "4x3":
        height = Inches(15)
    else:
        height = Inches(28.58 / 2.54)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if ratio == "16x9":
        tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    if '<underline>' in v:
        p = tf.paragraphs[0]
        add_underlined_text(p, v)
    else:
        tf.text = v
    p = tf.paragraphs[0]
    p.font.size = Pt(60)
    p.alignment = PP_ALIGN.CENTER


def write_pptx(psalm, ratio, colour, underline, output_folder):
    tmpfile = tempfile.SpooledTemporaryFile()
    # set up pptx:
    prs = Presentation(
        os.path.join(
            "..",
            "masters",
            ratio + "_" + colour + ".pptx"
        )
    )
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    # title slide:
    title.text = psalm['name']
    if psalm['book'] == 'Sing Psalms':
        subtitle.text = psalm['metre'] + u'\n\u00A9 Free Church of Scotland'
    else:
        subtitle.text = psalm['metre']
    # stanzas
    for v in psalm['stanzas']:
        if not underline:
            v = remove_markup(v)
        add_slide(prs, v, ratio, underline)
    # save tmp pptx
    prs.save(tmpfile)
    # superscript verse numbers
    with warnings.catch_warnings():
        warnings.simplefilter("ignore")
        superscript_pptx(tmpfile)
    # reload and save again to fix powerpoint complaining
    prs_ = Presentation(tmpfile)
    prs_.save(os.path.join(output_folder, psalm['file_name'] + '.pptx'))

    tmpfile.close()


def sing_psalms(ratio="16x9", colour="b_w", underline=False):
    # sing psalms
    folder_ids = [ratio, colour]
    if underline:
        folder_ids.append("underlined")
    file_name = "Sing Psalms"
    output_folder = make_output_folder(["PowerPoint", '_'.join(folder_ids), file_name])
    psalms = load_sing_psalms()
    for psalm in psalms:
        write_pptx(psalm, ratio, colour, underline, output_folder)

    pdf_folder = convert2pdf(output_folder)
    return output_folder, pdf_folder


def trad_psalms(ratio="16x9", colour="b_w", underline=False):
    # trad psalms
    folder_ids = [ratio, colour]
    if underline:
        folder_ids.append("underlined")
    file_name = "Scottish Psalter"
    output_folder = make_output_folder(["PowerPoint", '_'.join(folder_ids), file_name])
    psalms = load_scottish_psalter()
    for psalm in psalms:
        write_pptx(psalm, ratio, colour, underline, output_folder)

    pdf_folder = convert2pdf(output_folder)
    return output_folder, pdf_folder


def convert2pptx(ratio="16x9", colour="b_w", underline=False):
    """Convert Psalms to pptx files"""
    out_folder, pdf_folder = sing_psalms(ratio=ratio, colour=colour, underline=underline)
    trad_psalms(ratio=ratio, colour=colour, underline=underline)

    out_folder = os.path.dirname(out_folder)
    pdf_folder = os.path.dirname(pdf_folder)

    # clean up:
    for folder in [out_folder, pdf_folder]:
        zip_folder(folder)
        remove_folder(folder)

if __name__ == '__main__':
    convert2pptx()
